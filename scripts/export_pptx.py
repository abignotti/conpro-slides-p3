"""
Export provisional del deck a PPTX (imagen-por-slide).

⚠️ Provisorio: captura cada slide renderizado y lo pega como imagen full-bleed,
así que NO es editable en PowerPoint. Se reemplazará por un export nativo-editable
(cajas de texto/tablas reales) — ver hito 6 en docs/project-spec.md.

Requisitos: playwright (con chromium) + python-pptx, y el server local corriendo:
    python3 -m http.server 8753
"""
import os
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FR = os.path.join(ROOT, "build", "pptx_frames")
os.makedirs(FR, exist_ok=True)
URL = "http://localhost:8753/index.html"

# 1) Capturar las slides a 1920x1080 (ocultando el chrome de previsualización)
with sync_playwright() as p:
    b = p.chromium.launch(headless=True)
    pg = b.new_page(viewport={"width": 1920, "height": 1080}, device_scale_factor=2)
    pg.goto(URL); pg.wait_for_load_state("networkidle"); pg.wait_for_timeout(1200)
    for sel in (".theme-picker", ".reveal .controls", ".reveal .progress"):
        pg.evaluate(f"var e=document.querySelector('{sel}'); if(e) e.style.display='none';")
    n = pg.evaluate("Reveal.getTotalSlides()")
    for i in range(n):
        pg.evaluate(f"Reveal.slide({i})"); pg.wait_for_timeout(650)
        pg.screenshot(path=f"{FR}/s{i+1:02d}.png")
    b.close()

# 2) Armar el PPTX 16:9 (13.333 x 7.5 in), 1 imagen full-bleed por slide
prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]
for i in range(n):
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(f"{FR}/s{i+1:02d}.png", 0, 0, width=prs.slide_width, height=prs.slide_height)
out = os.path.join(ROOT, "Escalamiento-Conpro.pptx")
prs.save(out)
print("PPTX:", out, "slides:", n, "size:", os.path.getsize(out))
