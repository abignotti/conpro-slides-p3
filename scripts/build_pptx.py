"""
Genera un PPTX NATIVO y EDITABLE de "Escalamiento de Conpro", lo más parecido
posible al deck web (Mostaza claro · Corporativo · chrome completo).

- Mapea el lienzo web 1920×1080 a un slide 13.333×7.5in: 1 px = 6350 EMU,
  tamaños de fuente = px × 0.5 pt. Así las posiciones replican el layout web.
- Texto, tablas, tarjetas y formas son NATIVOS (editables en PowerPoint).
- Los gráficos se insertan como imagen capturada del deck web (fiel), porque
  recrear cada gráfico nativo no aporta a la edición de las compañeras.

Requisitos: python-pptx + playwright (chromium) + el server local:
    python3 -m http.server 8753
Uso: python3 scripts/build_pptx.py
Salida: Escalamiento-Conpro.pptx (en la raíz)
"""
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
URL = "http://localhost:8753/index.html"
CHART_DIR = os.path.join(ROOT, "build", "pptx_charts")
os.makedirs(CHART_DIR, exist_ok=True)

# ---- mapeo px (lienzo 1920×1080) -> unidades pptx ----
EMU_PER_PX = 6350
def X(px): return Emu(int(round(px * EMU_PER_PX)))
def FS(px): return Pt(px * 0.5)

# ---- colores Mostaza claro ----
C = {
    "bg": "FBFAF7", "bg_alt": "F2EEE4", "text": "1A1A17", "muted": "76726A",
    "accent": "FCCE4E", "on_accent": "1A1A17", "grid": "E4DECF",
}
def rgb(k): return RGBColor.from_string(C[k])

TITLE_FONT = "Montserrat"   # display (Corporativo)
BODY_FONT = "Inter"         # cuerpo

# escala (px) — desde tokens.css
HERO, H1, H2, STAT, BODY, CAP, KICK = 168, 72, 40, 66, 26, 24, 24

PAD = 80
RIGHT = 1920 - PAD  # 1840

# ============================================================ helpers
def set_bg(slide, key):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(key)

def _set_run(r, text, size_px, color="text", bold=False, font=BODY_FONT, caps=False):
    r.text = text.upper() if caps else text
    r.font.size = FS(size_px)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = rgb(color)

def text(slide, x, y, w, h, runs, align="l", anchor="t", line=1.12, wrap=True):
    """runs: lista de párrafos; cada párrafo = lista de dicts de run."""
    tb = slide.shapes.add_textbox(X(x), X(y), X(w), X(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    for m in ("left", "right", "top", "bottom"):
        setattr(tf, f"margin_{m}", 0)
    tf.vertical_anchor = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}[anchor]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}[align]
        p.line_spacing = line
        for rd in para:
            _set_run(p.add_run(), **rd)
    return tb

def R(t, sz, color="text", bold=False, font=BODY_FONT, caps=False):
    return {"text": t, "size_px": sz, "color": color, "bold": bold, "font": font, "caps": caps}

def rect(slide, x, y, w, h, fill=None, line=None, round_=False, line_w=1):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        X(x), X(y), X(w), X(h))
    if round_:
        try: shp.adjustments[0] = 0.06
        except Exception: pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = rgb(line); shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

def hairline(slide, x, y, w):
    rect(slide, x, y, w, 1.5, fill="grid")

def header(slide, kicker, folio):
    text(slide, PAD, 46, 1200, 40, [[R(kicker, KICK, "muted", bold=True, font=TITLE_FONT, caps=True)]])
    if folio:
        text(slide, RIGHT - 1200, 46, 1200, 40, [[R(folio, CAP, "muted")]], align="r")
    hairline(slide, PAD, 96, RIGHT - PAD)

def footer(slide):
    hairline(slide, PAD, 1000, RIGHT - PAD)
    text(slide, PAD, 1012, 1400, 40, [[R("Escalamiento de Conpro", CAP, "muted")]])

def chart_img(slide, cid, x, y, w, h):
    img = os.path.join(CHART_DIR, cid + ".png")
    if os.path.exists(img):
        slide.shapes.add_picture(img, X(x), X(y), width=X(w))
    else:
        rect(slide, x, y, w, h, fill="bg_alt", line="grid")

# ============================================================ renderers
def s_portada(slide, d):
    set_bg(slide, "bg_alt")
    text(slide, PAD, 46, 900, 40, [[R("Conpro · 2026", CAP, "muted", bold=True, font=TITLE_FONT, caps=True)]])
    text(slide, RIGHT - 1200, 46, 1200, 40, [[R("ICS3913 · Evaluación de Proyectos", CAP, "muted", bold=True, font=TITLE_FONT, caps=True)]], align="r")
    text(slide, PAD, 170, 1400, 40, [[R("Proyecto de título · Informe 3", KICK, "muted", bold=True, font=TITLE_FONT, caps=True)]])
    text(slide, PAD, 215, 1500, 320, [[R("Escalamiento", 128, "text", bold=True, font=TITLE_FONT)],
                                        [R("de Conpro", 128, "text", bold=True, font=TITLE_FONT)]], line=0.98)
    rect(slide, PAD, 560, 200, 6, fill="accent", round_=True)
    text(slide, PAD, 600, 1000, 110, [[R("Evaluación financiera y estratégica del escalamiento desde la operación informal por WhatsApp hacia una plataforma digital.", 30, "muted")]], line=1.45)
    hairline(slide, PAD, 952, RIGHT - PAD)
    text(slide, PAD, 964, 1500, 90, [
        [R("Grupo 06", CAP, "text", bold=True)],
        [R("Agustín Bignotti · María Soledad Bravo · Javiera Hellwig · Barbara Küllmer · Emilia Luttecke · Dafne Valdivia", CAP, "muted")],
    ], line=1.3)
    text(slide, RIGHT - 600, 970, 600, 40, [[R("Junio 2026", CAP, "muted", caps=True)]], align="r")

def s_divisor(slide, d):
    set_bg(slide, "bg_alt")
    header(slide, "Evaluación de Proyectos", d["folio"])
    text(slide, PAD, 240, 520, 460, [[R(d["num"], 300, "accent", bold=True, font=TITLE_FONT)]], anchor="m", line=0.8)
    lx = 560
    text(slide, lx, 250, 1200, 40, [[R("Sección " + d["num"], KICK, "muted", bold=True, font=TITLE_FONT, caps=True)]])
    text(slide, lx, 290, 1280, 110, [[R(d["title"], 72, "text", bold=True, font=TITLE_FONT)]], line=1.0)
    secs = ["El cliente y el negocio", "Caso Base Optimizado", "El proyecto de escalamiento", "Modelo financiero", "Robustez y decisión"]
    y = 430
    for i, s in enumerate(secs, 1):
        on = (i == d["active"])
        rect(slide, lx, y, 46, 46, fill=("accent" if on else None), line=(None if on else "grid"), round_=True)
        text(slide, lx, y, 46, 46, [[R(f"{i:02d}", 21, "on_accent" if on else "muted", bold=True)]], align="c", anchor="m")
        text(slide, lx + 66, y, 1000, 46, [[R(s, 27, "text" if on else "muted", bold=on)]], anchor="m")
        y += 62
    footer(slide)

def s_cita(slide, d):
    set_bg(slide, "bg_alt")
    header(slide, "Evaluación de Proyectos", d["folio"])
    rect(slide, PAD, 360, 6, 360, fill="accent", round_=True)
    sz = d.get("size", 64)
    text(slide, 150, 330, 1500, 420, [[R(d["text"], sz, "text", bold=True, font=TITLE_FONT)]], anchor="m", line=1.16)
    text(slide, 150, 760, 1500, 40, [[R(d["attr"], BODY, "muted", caps=True)]])
    footer(slide)

def s_bullets_marker(slide, d):
    """split imagen / split chart / split hero comparten columna izquierda con bullets de marcador."""
    pass

def _kicker_title(slide, kicker, title, ty=210, th=150, tw=1700):
    text(slide, PAD, 150, 1400, 40, [[R(kicker, KICK, "muted", bold=True, font=TITLE_FONT, caps=True)]])
    text(slide, PAD, ty, tw, th, [[R(title, H1, "text", bold=True, font=TITLE_FONT)]], line=1.04)

def _marker_bullets(slide, items, x, y, w, gap=22, sz=BODY):
    for it in items:
        rect(slide, x, y + 9, 16, 16, fill="accent", round_=True)
        n = text(slide, x + 34, y, w - 34, 80, [[R(it, sz, "text")]], line=1.45)
        # estimar alto por nº de líneas (~ chars/linea)
        import math
        lines = max(1, math.ceil(len(it) / 42))
        y += lines * 38 + gap
    return y

def s_split_image(slide, d):
    set_bg(slide, "bg")
    header(slide, "Evaluación de Proyectos", d["folio"])
    _kicker_title(slide, d["kicker"], d["title"], ty=210, th=120, tw=820)
    text(slide, PAD, 360, 820, 120, [[R(d["intro"], BODY, "muted")]], line=1.5)
    _marker_bullets(slide, d["bullets"], PAD, 510, 820)
    rect(slide, 1000, 250, 840, 600, fill="bg_alt", line="grid", round_=True)
    text(slide, 1000, 520, 840, 60, [[R(d["image"], CAP, "muted")]], align="c", anchor="m")
    footer(slide)

def s_split_chart(slide, d):
    set_bg(slide, "bg")
    header(slide, "Evaluación de Proyectos", d["folio"])
    _kicker_title(slide, d["kicker"], d["title"], ty=200, th=210, tw=860)
    y0 = 450
    if d.get("intro"):
        text(slide, PAD, 420, 820, 60, [[R(d["intro"], BODY, "muted")]], line=1.5)
        y0 = 490
    _marker_bullets(slide, d["bullets"], PAD, y0, 840)
    # tarjeta + gráfico
    cx, cy, cw, chh = 980, 250, 860, 600
    rect(slide, cx, cy, cw, chh, fill="bg_alt", line="grid", round_=True)
    text(slide, cx + 44, cy + 36, cw - 88, 60, [[R(d["cardtitle"], CAP, "muted", bold=True, font=TITLE_FONT, caps=True)]], line=1.2)
    chart_img(slide, d["chart"], cx + 44, cy + 110, cw - 88, chh - 160)
    footer(slide)

def s_chart_full(slide, d):
    set_bg(slide, "bg")
    header(slide, "Evaluación de Proyectos", d["folio"])
    _kicker_title(slide, d["kicker"], d["title"], ty=200, th=160)
    chart_img(slide, d["chart"], PAD, 400, RIGHT - PAD, 520)
    footer(slide)

def s_split_hero(slide, d):
    set_bg(slide, "bg")
    header(slide, "Evaluación de Proyectos", d["folio"])
    _kicker_title(slide, d["kicker"], d["title"], ty=200, th=210, tw=820)
    _marker_bullets(slide, d["bullets"], PAD, 510, 820)
    cx, cy, cw, chh = 980, 280, 860, 470
    rect(slide, cx, cy, cw, chh, fill="bg_alt", line="grid", round_=True)
    text(slide, cx + 56, cy + 56, cw - 112, 40, [[R(d["herolabel"], CAP, "muted", bold=True, font=TITLE_FONT, caps=True)]])
    text(slide, cx + 56, cy + 110, cw - 112, 130, [[R(d["hero"], 116, "text", bold=True, font=TITLE_FONT)]], line=0.95, wrap=False)
    text(slide, cx + 56, cy + 290, cw - 112, 150, [[R(d["herosub"], BODY, "muted")]], line=1.45)
    footer(slide)

def s_stat_grid(slide, d):
    set_bg(slide, "bg")
    header(slide, "Evaluación de Proyectos", d["folio"])
    _kicker_title(slide, d["kicker"], d["title"], ty=200, th=160)
    n = len(d["cells"]); gap = 1; total = RIGHT - PAD
    cw = (total - gap * (n - 1)) / n
    y, h = 470, 300
    hairline(slide, PAD, y - 1, total)
    for i, (big, lab) in enumerate(d["cells"]):
        x = PAD + i * (cw + gap)
        on = (i == d["accent"])
        if on: rect(slide, x, y, cw, h, fill="accent")
        text(slide, x + 44, y + 48, cw - 88, 90, [[R(big, STAT, "on_accent" if on else "text", bold=True, font=TITLE_FONT)]])
        text(slide, x + 44, y + 150, cw - 88, 120, [[R(lab, BODY, "on_accent" if on else "muted")]], line=1.4)
    hairline(slide, PAD, y + h, total)
    footer(slide)

def s_kpi(slide, d):
    set_bg(slide, "bg")
    header(slide, "Evaluación de Proyectos", d["folio"])
    _kicker_title(slide, d["kicker"], d["title"], ty=200, th=120)
    n = len(d["cards"]); gap = 24; total = RIGHT - PAD
    cw = (total - gap * (n - 1)) / n
    y, h = 440, 300
    for i, (lab, big, sub) in enumerate(d["cards"]):
        x = PAD + i * (cw + gap)
        on = (i == d["accent"])
        rect(slide, x, y, cw, h, fill="accent" if on else "bg_alt", line=None if on else "grid", round_=True)
        text(slide, x + 36, y + 36, cw - 72, 70, [[R(lab, CAP, "on_accent" if on else "muted", bold=True, font=TITLE_FONT, caps=True)]], line=1.15)
        text(slide, x + 36, y + 150, cw - 72, 70, [[R(big, 56, "on_accent" if on else "text", bold=True, font=TITLE_FONT)]])
        text(slide, x + 36, y + 232, cw - 72, 50, [[R(sub, CAP, "on_accent" if on else "muted")]])
    footer(slide)

def s_timeline(slide, d):
    set_bg(slide, "bg")
    header(slide, "Evaluación de Proyectos", d["folio"])
    _kicker_title(slide, d["kicker"], d["title"], ty=200, th=120)
    steps = d["steps"]; n = len(steps); total = RIGHT - PAD
    colw = total / n
    y = 470
    rect(slide, PAD + 22, y + 21, total - 44, 1.5, fill="grid")
    for i, (t, p) in enumerate(steps):
        x = PAD + i * colw
        on = (i + 1 == d["active"])
        rect(slide, x, y, 44, 44, fill="accent" if on else "bg", line=None if on else "grid", round_=True)
        text(slide, x, y, 44, 44, [[R(f"{i+1:02d}", 24, "on_accent" if on else "muted", bold=True)]], align="c", anchor="m")
        text(slide, x, y + 70, colw - 24, 50, [[R(t, 28, "text", bold=True, font=TITLE_FONT)]], line=1.1)
        text(slide, x, y + 130, colw - 24, 160, [[R(p, 24, "muted")]], line=1.45)
    footer(slide)

def s_tabla(slide, d):
    set_bg(slide, "bg")
    header(slide, "Evaluación de Proyectos", d["folio"])
    _kicker_title(slide, d["kicker"], d["title"], ty=200, th=120)
    yp = 330
    if d.get("pills"):
        px_ = PAD
        for pill in d["pills"]:
            w = 60 + len(pill) * 12
            rect(slide, px_, yp, w, 52, fill="bg_alt", line="grid", round_=True)
            text(slide, px_, yp, w, 52, [[R(pill, 22, "text", bold=True)]], align="c", anchor="m")
            px_ += w + 16
        yp = 420
    _table(slide, d, PAD, yp, RIGHT - PAD)
    footer(slide)

def _table(slide, d, x, y, w):
    headers = d["headers"]; rows = d["rows"]; nc = len(headers)
    nr = len(rows) + 1
    rh = 86
    tbl = slide.shapes.add_table(nr, nc, X(x), X(y), X(w), X(rh * nr)).table
    tbl.first_row = False; tbl.horz_banding = False
    # anchos: primera col más ancha
    tbl.columns[0].width = X(w * 0.34)
    for c in range(1, nc):
        tbl.columns[c].width = X(w * 0.66 / (nc - 1))
    def style(cell, txt, sz, color, bold, align, fill):
        cell.fill.solid(); cell.fill.fore_color.rgb = rgb(fill)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = X(20); cell.margin_right = X(20)
        cell.margin_top = 0; cell.margin_bottom = 0
        p = cell.text_frame.paragraphs[0]
        p.alignment = {"l": PP_ALIGN.LEFT, "r": PP_ALIGN.RIGHT, "c": PP_ALIGN.CENTER}[align]
        _set_run(p.add_run(), txt, sz, color=color, bold=bold, font=BODY_FONT)
    for c, htxt in enumerate(headers):
        style(tbl.cell(0, c), htxt, CAP, "muted", True, "l" if c == 0 else "r", "bg")
    for r, row in enumerate(rows, 1):
        acc = (r - 1 == d.get("accent", -1))
        for c, val in enumerate(row):
            keycell = acc and c == d.get("accent_col", 1)
            fill = "accent" if keycell else ("bg_alt" if acc else "bg")
            color = "on_accent" if keycell else "text"
            style(tbl.cell(r, c), val, 28, color, (acc or keycell), "l" if c == 0 else "r", fill)

def s_tabla_grafico(slide, d):
    set_bg(slide, "bg")
    header(slide, "Evaluación de Proyectos", d["folio"])
    _kicker_title(slide, d["kicker"], d["title"], ty=200, th=120)
    _table(slide, d, PAD, 430, 820)
    chart_img(slide, d["chart"], 1000, 420, 840, 420)
    footer(slide)

def s_comparacion(slide, d):
    set_bg(slide, "bg")
    header(slide, "Evaluación de Proyectos", d["folio"])
    _kicker_title(slide, d["kicker"], d["title"], ty=200, th=120)
    total = RIGHT - PAD; gap = 32; cw = (total - gap) / 2; y = 440; h = 420
    for i, card in enumerate([d["left"], d["right"]]):
        x = PAD + i * (cw + gap)
        acc = card.get("accent")
        rect(slide, x, y, cw, h, fill=None, line="grid", round_=True)
        rect(slide, x, y, cw, 80, fill="accent" if acc else "bg_alt")
        text(slide, x + 40, y, cw - 380, 80, [[R(card["name"], 34, "on_accent" if acc else "text", bold=True, font=TITLE_FONT)]], anchor="m")
        text(slide, x + cw - 380, y, 340, 80, [[R(card["tag"], CAP, "on_accent" if acc else "muted", bold=True, caps=True)]], align="r", anchor="m", wrap=False)
        my = y + 110
        for lab, val in card["metrics"]:
            text(slide, x + 40, my, cw * 0.55, 50, [[R(lab, BODY, "muted")]], anchor="m")
            text(slide, x + cw - 360, my, 320, 50, [[R(val, 32, "text", bold=True)]], align="r", anchor="m")
            my += 100
    footer(slide)

def s_matriz(slide, d):
    set_bg(slide, "bg")
    header(slide, "Evaluación de Proyectos", d["folio"])
    _kicker_title(slide, d["kicker"], d["title"], ty=200, th=120)
    quads = d["quads"]; total = RIGHT - PAD; gap = 1; cw = (total - gap * 3) / 4
    y, h = 440, 360
    hairline(slide, PAD, y - 1, total)
    for i, (letter, lab, body) in enumerate(quads):
        x = PAD + i * (cw + gap); on = (i == 0)
        if on: rect(slide, x, y, cw, h, fill="accent")
        text(slide, x + 36, y + 36, cw - 72, 110, [[R(letter, 96, "on_accent" if on else "text", bold=True, font=TITLE_FONT)]], line=0.9)
        text(slide, x + 36, y + 150, cw - 72, 40, [[R(lab, CAP, "on_accent" if on else "muted", bold=True, caps=True)]])
        text(slide, x + 36, y + 200, cw - 72, 150, [[R(body, 24, "on_accent" if on else "muted")]], line=1.45)
    hairline(slide, PAD, y + h, total)
    footer(slide)

def s_bullets(slide, d):
    set_bg(slide, "bg")
    header(slide, "Evaluación de Proyectos", d["folio"])
    _kicker_title(slide, d["kicker"], d["title"], ty=190, th=120)
    y = 360
    for i, (h3, p) in enumerate(d["items"]):
        on = d.get("all_accent") or (i == 0)
        hairline(slide, PAD, y, RIGHT - PAD)
        rect(slide, PAD, y + 24, 54, 54, fill="accent" if on else "bg_alt", line=None if on else "grid", round_=True)
        text(slide, PAD, y + 24, 54, 54, [[R(f"{i+1:02d}", 26, "on_accent" if on else "muted", bold=True)]], align="c", anchor="m")
        text(slide, PAD + 90, y + 18, 1600, 50, [[R(h3, H2, "text", bold=True, font=TITLE_FONT)]])
        text(slide, PAD + 90, y + 78, 1600, 70, [[R(p, BODY, "muted")]], line=1.5)
        y += 200
    footer(slide)


# ============================================================ datos de las 28 slides
SECS = ["El cliente y el negocio", "Caso Base Optimizado", "El proyecto de escalamiento", "Modelo financiero", "Robustez y decisión"]
SLIDES = [
    ("portada", {}),
    ("divisor", {"folio": "02 / 28", "num": "01", "title": SECS[0], "active": 1}),
    ("split_image", {"folio": "03 / 28", "kicker": "El cliente", "title": "Juan, el cliente objetivo",
        "intro": "40 años, vive en Providencia y es papá de dos hijos. Un hogar de ingresos medios-altos que planifica sus compras y busca optimizar el gasto.",
        "bullets": ["Valora la calidad, el origen y la confianza de lo que lleva a su casa.",
                    "Pero lo premium es caro: elige entre pagar de más o resignar calidad."],
        "image": "Foto · Juan (cliente objetivo)"}),
    ("cita", {"folio": "04 / 28", "text": "Conpro agrega la demanda de muchos hogares para acceder a productos de calidad a un precio justo: el que cada uno, por separado, no alcanzaría.", "attr": "Propuesta de valor"}),
    ("split_chart", {"folio": "05 / 28", "kicker": "Análisis de mercado", "title": "Compite con las especializadas, no con el supermercado",
        "bullets": ["Súper y mayoristas: más baratos, pero de menor calidad (no son competencia).",
                    "Naturistas y especializadas: calidad similar, pero más caras.",
                    "Conpro: la misma calidad, a un precio más conveniente."],
        "chart": "chart-precios", "cardtitle": "Precio por producto · Conpro vs mercado"}),
    ("timeline", {"folio": "06 / 28", "kicker": "Cómo funciona hoy", "title": "El pedido colectivo, por WhatsApp", "active": 1,
        "steps": [("Disponibilidad", "Ricardo informa producto, condiciones y volumen mínimo de la semana."),
                  ("Inscripción", "Cada interesado se suma al pedido colectivo de su grupo."),
                  ("Volumen mínimo", "Si se alcanza el mínimo del proveedor, la compra se concreta."),
                  ("Pago", "Se consolidan las transferencias y se paga al proveedor."),
                  ("Entrega", "Despacho a domicilio o retiro en la conserjería de Ricardo.")]}),
    ("bullets", {"folio": "07 / 28", "kicker": "La dependencia de Ricardo", "title": "Toda la gestión pasa por Ricardo",
        "items": [("Pedidos", "Reúne la demanda y envía las órdenes de compra a cada proveedor."),
                  ("Pagos", "Recibe transferencias, revisa comprobantes, valida pagos y transfiere los fondos."),
                  ("Entrega", "Coordina el despacho y el retiro de los pedidos en su conserjería.")]}),
    ("stat_grid", {"folio": "08 / 28", "kicker": "Costo de oportunidad", "title": "El tiempo operativo rinde menos de lo que cuesta", "accent": 2,
        "cells": [("16 h/mes", "Dedicación de Ricardo a tareas administrativas"),
                  ("$5.500/h", "Lo que rinde hoy esa hora dedicada al negocio"),
                  ("$16.700/h", "Su costo de oportunidad en otras actividades")]}),
    ("divisor", {"folio": "09 / 28", "num": "02", "title": SECS[1], "active": 2}),
    ("split_hero", {"folio": "10 / 28", "kicker": "Caso Base Optimizado", "title": "Automatizar los pagos, sin cambiar el modelo",
        "bullets": ["Mantiene la comunidad de WhatsApp y la compra colectiva.",
                    "Automatiza los pagos con Fintoc: links de pago y registro automático.",
                    "Libera ~2,5 h/semana para difusión y captación."],
        "herolabel": "Inversión inicial", "hero": "$295.000", "herosub": "Formalización de la marca y términos & condiciones."}),
    ("kpi", {"folio": "11 / 28", "kicker": "Resultados del CBO", "title": "Más valor con una inversión mínima", "accent": 3,
        "cards": [("Ingreso mensual · año 5", "$126.850", "▲ 20% vs. hoy"),
                  ("Ingreso por hora", "$7.930", "desde $5.500"),
                  ("Inversión", "$295.000", "marca + Fintoc"),
                  ("VAN · 5 años", "$2,41 M", "payback al año 1")]}),
    ("cita", {"folio": "12 / 28", "text": "El CBO mejora la operación, pero no rompe el techo: el crecimiento sigue atado a Ricardo y a un punto de retiro limitado.", "attr": "El límite del Caso Base · techo ~180 hogares"}),
    ("divisor", {"folio": "13 / 28", "num": "03", "title": SECS[2], "active": 3}),
    ("timeline", {"folio": "14 / 28", "kicker": "El proyecto · plataforma B2C", "title": "La compra colectiva, ahora en una plataforma", "active": 5,
        "steps": [("Grupos por producto", "El usuario se une al grupo del producto que le interesa."),
                  ("Meta de volumen", "El grupo sigue activo hasta alcanzar el mínimo del proveedor."),
                  ("Pago y cierre", "Al llegar a la meta se habilita el pago y se cierra el grupo."),
                  ("Transferencia única", "Conpro consolida los pagos y transfiere al proveedor."),
                  ("Sin cuello de botella", "El flujo corre solo: escalable y menos dependiente de Ricardo.")]}),
    ("cita", {"folio": "15 / 28", "text": "¿Realmente vale la pena invertir en esta plataforma?", "attr": "La pregunta que responde el modelo financiero", "size": 80}),
    ("divisor", {"folio": "16 / 28", "num": "04", "title": SECS[3], "active": 4}),
    ("split_chart", {"folio": "17 / 28", "kicker": "Modelo financiero · la demanda", "title": "La variable más incierta del modelo",
        "intro": "Modelo logístico de difusión, con tres factores:",
        "bullets": ["Boca a boca: el crecimiento orgánico de la comunidad.",
                    "Marketing: inversión para captar nuevos hogares.",
                    "Tasa de fuga: cuántos usuarios se pierden."],
        "chart": "chart-demanda", "cardtitle": "Ingreso bruto anual · escenario base (M$)"}),
    ("stat_grid", {"folio": "18 / 28", "kicker": "Costos", "title": "Una inversión acotada, costos livianos", "accent": 0,
        "cells": [("$2.445.000", "Inversión inicial: plataforma ($1,6 M) y gastos legales"),
                  ("$310.000", "Costos fijos al año: patente, hosting y servidor"),
                  ("50%", "De las utilidades reinvertidas en marketing (etapa inicial)")]}),
    ("tabla_grafico", {"folio": "19 / 28", "kicker": "Ganancias proyectadas", "title": "La ganancia crece con la escala",
        "headers": ["Año", "Ganancia (M$)"], "accent": 4, "accent_col": 1,
        "rows": [["Año 1", "0,5"], ["Año 2", "1,6"], ["Año 3", "2,8"], ["Año 4", "3,9"], ["Año 5", "5,0"]],
        "chart": "chart-ganancias"}),
    ("tabla", {"folio": "20 / 28", "kicker": "Indicadores de rentabilidad", "title": "VAN positivo en los tres escenarios",
        "pills": ["WACC 29,59%", "Payback base · 3 años"],
        "headers": ["Escenario", "VAN", "TIR"], "accent": 1, "accent_col": 1,
        "rows": [["Pesimista", "$1.123.525", "45,7%"], ["Conservador", "$2.386.311", "57,9%"], ["Optimista", "$3.969.551", "69,3%"]]}),
    ("divisor", {"folio": "21 / 28", "num": "05", "title": SECS[4], "active": 5}),
    ("split_chart", {"folio": "22 / 28", "kicker": "Análisis de sensibilidad", "title": "Dos variables mandan el resultado",
        "bullets": ["El ingreso por hogar y la tasa de fuga explican casi todo el VAN.",
                    "No basta atraer usuarios: deben comprar seguido y mantenerse activos."],
        "chart": "chart-tornado", "cardtitle": "Impacto de cada variable en el VAN"}),
    ("split_chart", {"folio": "23 / 28", "kicker": "Análisis de márgenes", "title": "Café y vino: espacio para crecer",
        "bullets": ["Huevos y queso ya están en su límite competitivo.",
                    "Café, vino y aceite tienen holgura para subir margen sin perder competitividad."],
        "chart": "chart-margenes", "cardtitle": "Aumento del VAN por producto"}),
    ("chart_full", {"folio": "24 / 28", "kicker": "La opción de esperar", "title": "Esperar puede crear más valor que invertir hoy", "chart": "chart-esperar"}),
    ("comparacion", {"folio": "25 / 28", "kicker": "Proyecto vs Caso Base", "title": "Mismo valor, una fracción de la inversión",
        "left": {"name": "Proyecto (plataforma)", "tag": "Alternativa", "metrics": [("Inversión", "$2.445.000"), ("VAN", "$2,39 M"), ("Payback · riesgo", "3 años · alto")]},
        "right": {"name": "Caso Base Optimizado", "tag": "Recomendado", "accent": True, "metrics": [("Inversión", "$295.000"), ("VAN", "$2,41 M"), ("Payback · riesgo", "1 año · bajo")]}}),
    ("matriz", {"folio": "26 / 28", "kicker": "Análisis estratégico", "title": "FODA de Conpro",
        "quads": [("F", "Fortalezas", "Cercanía con el cliente y una comunidad preexistente con confianza establecida."),
                  ("O", "Oportunidades", "Crecimiento del e-commerce de alimentos y un contexto que impulsa el ahorro."),
                  ("D", "Debilidades", "Dependencia de proveedores y del volumen mínimo; sin marca ni reputación digital."),
                  ("A", "Amenazas", "Presión competitiva multicanal, imitación del modelo y riesgo de desintermediación.")]}),
    ("bullets", {"folio": "27 / 28", "kicker": "Recomendaciones", "title": "Optimizar hoy, escalar cuando se justifique", "all_accent": True,
        "items": [("No invertir todavía", "El VAN del proyecto es positivo, pero menor que el del CBO: hoy no es la alternativa más conveniente."),
                  ("Operar el CBO", "Automatizar pagos, validar la demanda y reducir la incertidumbre con baja inversión."),
                  ("Definir un gatillo", "Escalar a la plataforma cuando se alcancen los hogares activos que la hacen viable.")]}),
    ("cita", {"folio": "28 / 28", "text": "El mayor valor hoy está en perfeccionar lo que ya funciona. La plataforma no se descarta: se posterga hasta que la demanda la justifique.", "attr": "Recomendación final · ¡Gracias!"}),
]

RENDER = {
    "portada": s_portada, "divisor": s_divisor, "cita": s_cita, "split_image": s_split_image,
    "split_chart": s_split_chart, "split_hero": s_split_hero, "chart_full": s_chart_full,
    "stat_grid": s_stat_grid, "kpi": s_kpi, "timeline": s_timeline, "tabla": s_tabla,
    "tabla_grafico": s_tabla_grafico, "comparacion": s_comparacion, "matriz": s_matriz, "bullets": s_bullets,
}

# ============================================================ captura de gráficos
def capture_charts():
    from playwright.sync_api import sync_playwright
    ids = ["chart-precios", "chart-demanda", "chart-ganancias", "chart-tornado", "chart-margenes", "chart-esperar"]
    with sync_playwright() as p:
        b = p.chromium.launch(headless=True)
        pg = b.new_page(viewport={"width": 1920, "height": 1080}, device_scale_factor=2)
        pg.goto(URL); pg.wait_for_load_state("networkidle"); pg.wait_for_timeout(1500)
        for sel in (".theme-picker", ".reveal .controls", ".reveal .progress"):
            pg.evaluate(f"var e=document.querySelector('{sel}'); if(e) e.style.display='none';")
        total = pg.evaluate("Reveal.getTotalSlides()")
        for i in range(total):
            pg.evaluate(f"Reveal.slide({i})"); pg.wait_for_timeout(700)
            for cid in ids:
                el = pg.query_selector(f"section.present canvas#{cid}")
                if el:
                    el.screenshot(path=os.path.join(CHART_DIR, cid + ".png"))
        b.close()

# ============================================================ build
def main():
    capture_charts()
    prs = Presentation()
    prs.slide_width = Emu(12192000)   # 13.333 in
    prs.slide_height = Emu(6858000)   # 7.5 in
    blank = prs.slide_layouts[6]
    for kind, d in SLIDES:
        slide = prs.slides.add_slide(blank)
        RENDER[kind](slide, d)
    out = os.path.join(ROOT, "Escalamiento-Conpro.pptx")
    prs.save(out)
    print("PPTX:", out, os.path.getsize(out), "bytes,", len(SLIDES), "slides")

if __name__ == "__main__":
    main()
